"""
設計書解析モジュール
各種フォーマットの設計書ファイルからテキストを抽出します。
対応形式: PDF, Excel (.xlsx/.xls), Word (.docx), PowerPoint (.pptx), Markdown (.md), テキスト (.txt)
"""

import io
from pathlib import Path
from typing import Any


def parse_document(uploaded_file: Any) -> str:
    """
    アップロードされたファイルを解析してテキスト文字列として返す。
    ファイル名の拡張子をもとに適切なパーサーを選択する。
    """
    file_name = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()

    if file_name.endswith(".pdf"):
        return _parse_pdf(file_bytes)
    elif file_name.endswith((".xlsx", ".xls")):
        return _parse_excel(file_bytes)
    elif file_name.endswith(".docx"):
        return _parse_word(file_bytes)
    elif file_name.endswith(".pptx"):
        return _parse_powerpoint(file_bytes)
    elif file_name.endswith((".md", ".txt")):
        return file_bytes.decode("utf-8", errors="ignore")
    else:
        # 不明な拡張子はテキストとして解析を試みる
        try:
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            return ""


def _parse_pdf(file_bytes: bytes) -> str:
    """PDFファイルからテキストとテーブルを抽出する"""
    try:
        import pdfplumber

        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(f"--- ページ {page_num} ---\n{text.strip()}")

                # テーブルも抽出
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        table_text = _format_table(table)
                        if table_text:
                            text_parts.append(f"[テーブル - ページ {page_num}]\n{table_text}")

        return "\n\n".join(text_parts)
    except ImportError:
        return "エラー: pdfplumberがインストールされていません。pip install pdfplumber を実行してください。"
    except Exception as e:
        return f"PDFの解析に失敗しました: {str(e)}"


def _parse_excel(file_bytes: bytes) -> str:
    """Excelファイルの全シートからテキストを抽出する"""
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== シート: {sheet_name} ===")

            rows = []
            for row in sheet.iter_rows(values_only=True):
                # 完全に空の行はスキップ
                if any(cell is not None and str(cell).strip() != "" for cell in row):
                    rows.append([str(cell) if cell is not None else "" for cell in row])

            if rows:
                text_parts.append(_format_table(rows))

        return "\n\n".join(text_parts)
    except Exception as e:
        return f"Excelの解析に失敗しました: {str(e)}"


def _parse_word(file_bytes: bytes) -> str:
    """Wordファイルから段落とテーブルのテキストを抽出する"""
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                text_parts.append("\n[テーブル]\n" + _format_table(rows))

        return "\n".join(text_parts)
    except Exception as e:
        return f"Wordファイルの解析に失敗しました: {str(e)}"


def _parse_powerpoint(file_bytes: bytes) -> str:
    """PowerPointファイルから各スライドのテキストを抽出する"""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"--- スライド {slide_num} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(text_parts)
    except Exception as e:
        return f"PowerPointファイルの解析に失敗しました: {str(e)}"


def _format_table(table_data: list) -> str:
    """リスト形式のテーブルデータを読みやすいテキスト形式に変換する"""
    if not table_data:
        return ""

    lines = []
    for row in table_data:
        cells = [str(cell).strip() for cell in row]
        lines.append(" | ".join(cells))

    return "\n".join(lines)
